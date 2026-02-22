import os
import io
import json
import datetime
from typing import List, Optional, Dict

import joblib
import pandas as pd
import numpy as np
import faiss

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from pydantic import BaseModel
from sklearn.ensemble import RandomForestClassifier
from sklearn.model_selection import train_test_split
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from dotenv import load_dotenv

# ✅ New Gemini SDK (replaces deprecated google.generativeai)
from google import genai

# ----------------------------
# Env
# ----------------------------
load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-1.5-flash")  # start with flash (fast + usually available)

# Create Gemini client only if key exists (so server can still run without it)
gemini_client = genai.Client(api_key=GEMINI_API_KEY) if GEMINI_API_KEY else None

app = FastAPI(title="Campus AI Analytics & RAG API")

# ----------------------------
# Embedding Model
# ----------------------------
print("Loading Embedding Model...")
embed_model = SentenceTransformer("all-MiniLM-L6-v2")

MODEL_DIR = "models"
INDEX_DIR = "indices"
os.makedirs(MODEL_DIR, exist_ok=True)
os.makedirs(INDEX_DIR, exist_ok=True)

METADATA_PATH = os.path.join(MODEL_DIR, "metadata.json")

# FAISS Indices storage (in-memory)
# courseId -> {"index": faiss.IndexFlatL2(dim), "chunks": [chunk_info, ...]}
indices: Dict[str, Dict] = {}


# ----------------------------
# Helpers: metadata
# ----------------------------
def save_metadata(metadata: dict) -> None:
    with open(METADATA_PATH, "w") as f:
        json.dump(metadata, f)


def load_metadata() -> dict:
    if os.path.exists(METADATA_PATH):
        with open(METADATA_PATH, "r") as f:
            return json.load(f)
    return {"current_version": "v1", "models": {}}


# ----------------------------
# ML: Risk prediction schemas
# ----------------------------
class StudentData(BaseModel):
    student_id: str
    subject_id: str
    attendance_percentage: Optional[float] = 0
    average_assignment_score: Optional[float] = 0
    quiz_average: Optional[float] = 0
    gpa: Optional[float] = 0
    classes_missed: Optional[int] = 0
    late_submissions: Optional[int] = 0
    quiz_attempts: Optional[int] = 0
    face_violation_count: Optional[int] = 0
    payment_delay_days: Optional[int] = 0
    previous_exam_score: Optional[float] = 0
    exam_result: Optional[int] = None  # 1 FAIL, 0 PASS


class PredictionResponse(BaseModel):
    student_id: str
    risk_score: float
    risk_level: str
    reasons: List[str]
    version: Optional[str] = None


@app.post("/train")
async def train(data: List[StudentData]):
    if len(data) < 10:
        raise HTTPException(status_code=400, detail="Not enough data to train. Need at least 10 records.")

    df = pd.DataFrame([d.dict() for d in data])

    # Feature Engineering
    df["attendance_risk"] = (df["attendance_percentage"] < 75).astype(int)
    df["low_gpa"] = (df["gpa"] < 2.5).astype(int)

    X = df.drop(columns=["student_id", "subject_id", "exam_result"])
    y = df["exam_result"]

    if y.isnull().any():
        raise HTTPException(status_code=400, detail="Target variable 'exam_result' is missing in some records.")

    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

    model = RandomForestClassifier(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)

    accuracy = model.score(X_test, y_test)

    metadata = load_metadata()
    version_num = len(metadata["models"]) + 1
    version = f"v{version_num}"
    model_filename = f"model_{version}.pkl"
    model_path = os.path.join(MODEL_DIR, model_filename)

    joblib.dump(model, model_path)

    metadata["current_version"] = version
    metadata["models"][version] = {
        "path": model_path,
        "trained_at": datetime.datetime.now().isoformat(),
        "accuracy": float(accuracy),
        "records_used": len(data),
    }
    save_metadata(metadata)

    return {"message": "Model trained successfully", "version": version, "accuracy": float(accuracy), "records_used": len(data)}


@app.post("/predict", response_model=PredictionResponse)
async def predict(data: StudentData):
    metadata = load_metadata()
    current_version = metadata.get("current_version", "v1")
    model_info = metadata.get("models", {}).get(current_version)

    if not model_info or not os.path.exists(model_info["path"]):
        res = fallback_predict(data)
        res["version"] = "fallback"
        return res

    model = joblib.load(model_info["path"])

    df = pd.DataFrame([data.dict()])
    df["attendance_risk"] = (df["attendance_percentage"] < 75).astype(int)
    df["low_gpa"] = (df["gpa"] < 2.5).astype(int)

    X = df.drop(columns=["student_id", "subject_id", "exam_result"])

    risk_score = model.predict_proba(X)[0][1]  # prob of FAIL (class 1)

    if risk_score >= 0.7:
        risk_level = "HIGH"
    elif risk_score >= 0.4:
        risk_level = "MEDIUM"
    else:
        risk_level = "LOW"

    reasons = []
    if data.attendance_percentage < 75:
        reasons.append("Low attendance")
    if data.quiz_average < 50:
        reasons.append("Low quiz performance")
    if data.face_violation_count > 3:
        reasons.append("Multiple face violations")
    if data.average_assignment_score < 50:
        reasons.append("Low assignment scores")

    return {
        "student_id": data.student_id,
        "risk_score": float(risk_score),
        "risk_level": risk_level,
        "reasons": reasons,
        "version": current_version,
    }


def fallback_predict(data: StudentData):
    score = 0.0
    reasons = []

    if data.attendance_percentage < 75:
        score += 0.4
        reasons.append("Low attendance")
    if data.quiz_average < 50:
        score += 0.2
        reasons.append("Low quiz performance")
    if data.average_assignment_score < 50:
        score += 0.2
        reasons.append("Low assignment scores")
    if data.face_violation_count > 3:
        score += 0.1
        reasons.append("Multiple face violations")
    if data.payment_delay_days > 15:
        score += 0.1
        reasons.append("Payment delays")

    score = min(score, 1.0)

    if score >= 0.7:
        risk_level = "HIGH"
    elif score >= 0.4:
        risk_level = "MEDIUM"
    else:
        risk_level = "LOW"

    return {"student_id": data.student_id, "risk_score": score, "risk_level": risk_level, "reasons": reasons}


# ----------------------------
# RAG: Schemas
# ----------------------------
class ChatRequest(BaseModel):
    courseId: str
    question: str
    top_k: Optional[int] = 5


class ChunkData(BaseModel):
    content: str
    metadata: Dict


def extract_text_from_file(file_content: bytes, filename: str) -> List[ChunkData]:
    chunks: List[ChunkData] = []
    ext = filename.split(".")[-1].lower()

    if ext == "pdf":
        reader = PdfReader(io.BytesIO(file_content))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(ChunkData(content=text.strip(), metadata={"page": i + 1}))

    elif ext == "docx":
        doc = Document(io.BytesIO(file_content))
        current_text = ""
        for i, para in enumerate(doc.paragraphs):
            current_text += (para.text or "") + "\n"
            if len(current_text) > 1000:
                chunks.append(ChunkData(content=current_text.strip(), metadata={"paragraph_index": i}))
                current_text = ""
        if current_text.strip():
            chunks.append(ChunkData(content=current_text.strip(), metadata={"paragraph_index": len(doc.paragraphs)}))

    elif ext == "pptx":
        prs = Presentation(io.BytesIO(file_content))
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += (shape.text or "") + " "
            if text.strip():
                chunks.append(ChunkData(content=text.strip(), metadata={"slide": i + 1}))

    else:
        # Generic text
        text = file_content.decode("utf-8", errors="ignore")
        for i in range(0, len(text), 1000):
            piece = text[i : i + 1000].strip()
            if piece:
                chunks.append(ChunkData(content=piece, metadata={"offset": i}))

    return chunks


@app.post("/process-material")
async def process_material(courseId: str = Form(...), materialId: str = Form(...), file: UploadFile = File(...)):
    content = await file.read()
    chunks = extract_text_from_file(content, file.filename)

    if not chunks:
        return {"message": "No text extracted", "chunks": []}

    texts = [c.content for c in chunks]
    embeddings = embed_model.encode(texts)

    dim = embeddings.shape[1]
    if courseId not in indices:
        indices[courseId] = {"index": faiss.IndexFlatL2(dim), "chunks": []}

    start_idx = len(indices[courseId]["chunks"])
    indices[courseId]["index"].add(embeddings.astype("float32"))

    processed_chunks = []
    for i, c in enumerate(chunks):
        chunk_info = {
            "id": f"{materialId}_{start_idx + i}",
            "materialId": materialId,
            "content": c.content,
            "metadata": c.metadata,
            "global_idx": start_idx + i,
        }
        indices[courseId]["chunks"].append(chunk_info)
        processed_chunks.append(chunk_info)

    return {"message": f"Processed {len(chunks)} chunks", "chunks": processed_chunks}


def build_prompt(question: str, chunks: list) -> str:
    sources = []
    for i, c in enumerate(chunks, start=1):
        meta = c.get("metadata", {})
        # Include metadata so Gemini can cite properly
        sources.append(f"[S{i}] meta={meta}\n{c['content']}")

    return f"""
You are a university lecture assistant.

RULES:
- Answer ONLY using the SOURCES.
- If the answer is not found, say exactly: "I can't find this in the uploaded lecture materials."
- Keep it student-friendly.
- Add citations like [S1], [S2] after each key statement.

QUESTION:
{question}

SOURCES:
{chr(10).join(sources)}
""".strip()


@app.post("/chat")
async def chat(request: ChatRequest):
    if request.courseId not in indices or not indices[request.courseId]["chunks"]:
        return {
            "answer": "I don't have any materials for this course yet. Please upload some lecture notes first.",
            "citations": [],
        }

    # 1) Embed question
    q_emb = embed_model.encode([request.question])

    # 2) Search FAISS
    D, I = indices[request.courseId]["index"].search(q_emb.astype("float32"), request.top_k)

    retrieved_chunks = []
    for idx in I[0]:
        if idx != -1 and idx < len(indices[request.courseId]["chunks"]):
            retrieved_chunks.append(indices[request.courseId]["chunks"][idx])

    if not retrieved_chunks:
        return {"answer": "I can't find this in the uploaded lecture materials.", "citations": []}

    prompt = build_prompt(request.question, retrieved_chunks)

    # 3) Call Gemini (new SDK)
    if not gemini_client:
        answer = "Gemini is not configured. Please set GEMINI_API_KEY (or GOOGLE_API_KEY) in your .env and restart."
    else:
        try:
            resp = gemini_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
            )
            answer = resp.text or "I couldn't generate an answer. Please try again."
        except Exception as e:
            print(f"Error calling Gemini: {e}")
            answer = "I'm sorry, I'm having trouble connecting to Gemini right now."

    return {
        "answer": answer,
        "citations": [
            {
                "source": f"S{i+1}",
                "materialId": c["materialId"],
                "metadata": c["metadata"],
                "snippet": c["content"][:150],
            }
            for i, c in enumerate(retrieved_chunks)
        ],
    }


if __name__ == "__main__":
    import uvicorn
    import sys

    port = 8000
    if len(sys.argv) > 1:
        try:
            port = int(sys.argv[1])
        except ValueError:
            pass

    uvicorn.run(app, host="0.0.0.0", port=port)
